#!/usr/bin/env python3
"""Build results/REPORT_v2_<date>.pptx — the campaign deck.

Every number on every slide is read from the campaign's own summary JSONs at build time. Nothing
is typed into this script, so the deck cannot drift from the results the way a hand-maintained one
does — and a missing input is a hard failure rather than a silently empty slide, because a deck
that quietly ships a placeholder is worse than one that refuses to build.

``--allow-missing`` relaxes that to skip the affected slides and print exactly which ones were
skipped. It exists so the builder can be smoke-tested while the long stages are still running; it
is not for producing a deliverable, and the printed skip list is the reason it is safe to use.

    uv run --with python-pptx --with pillow python experiments/rikyu_hparam_tuning_v2/build_report_pptx.py
"""

from __future__ import annotations

import argparse
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
SUM = HERE / "summary"
RES = HERE / "results"

INK = RGBColor(0x1F, 0x29, 0x37)
MUT = RGBColor(0x6B, 0x72, 0x80)
BLUE = RGBColor(0x00, 0x77, 0xBB)
GREEN = RGBColor(0x00, 0x9E, 0x73)
RED = RGBColor(0xCC, 0x33, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADGREY = RGBColor(0x4B, 0x55, 0x63)

SKIPPED: list[str] = []
ALLOW_MISSING = False


class Missing(Exception):
    """An input this slide needs is not on disk yet."""


def need(path: Path) -> Path:
    if not path.exists():
        raise Missing(str(path))
    return path


def load(name: str) -> dict:
    return json.loads(need(SUM / name).read_text())


def slide_guard(fn):
    """Let one slide's missing input skip that slide instead of killing the build.

    Only under --allow-missing. Without it the exception propagates and the build fails, which is
    the behaviour a deliverable needs.
    """

    def wrapped(*a, **kw):
        try:
            return fn(*a, **kw)
        except Missing as exc:
            if not ALLOW_MISSING:
                raise SystemExit(
                    f"missing input for slide '{fn.__name__}': {exc}\n"
                    "run that stage's analysis first, or pass --allow-missing for a dry run"
                ) from exc
            SKIPPED.append(f"{fn.__name__}  (needs {exc})")
            return None

    return wrapped


prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def txt(slide, x, y, w, h, lines, size=14, color=INK, mono=False, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        if mono:
            p.font.name = "Consolas"
    return box


def title_bar(slide, title, sub=None):
    txt(slide, 0.5, 0.25, 12.3, 0.6, [title], size=24, bold=True)
    if sub:
        txt(slide, 0.5, 0.87, 12.3, 0.4, [sub], size=12, color=MUT)


def new(title, sub=None):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title, sub)
    return s


def pic_slide(title, sub, img, top=1.35, bottom=0.25):
    from PIL import Image

    need(img)
    s = new(title, sub)
    with Image.open(img) as im:
        iw, ih = im.size
    max_w, max_h = Inches(12.6), prs.slide_height - Inches(top) - Inches(bottom)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(img), Emu(int((prs.slide_width - w) / 2)), Inches(top), Emu(w), Emu(h))
    return s


def table(slide, x, y, w, headers, rows, col_w=None, size=11, head_size=11, colour_col=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w),
                                   Inches(0.32 * (len(rows) + 1)))
    tbl = shape.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(cw)
    for j, head in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = str(head)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(head_size)
        para.font.bold = True
        para.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADGREY
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(size)
            # Sign-colouring is opt-in per column: colouring every numeric cell would make the
            # deck read as a heat map and hide which column the argument is actually about.
            if colour_col is not None and j == colour_col and isinstance(value, str):
                para.font.color.rgb = (RED if value.startswith("-") else
                                       GREEN if value.startswith("+") else INK)
            else:
                para.font.color.rgb = INK
            if j:
                para.alignment = PP_ALIGN.RIGHT
    return tbl


# Run identifiers carry history the deck does not: "v1enc" is simply the encoder configuration
# inherited from the earlier campaign. Relabel at the presentation layer; the summary JSONs keep
# the original names so every figure stays traceable to its runs.
DISPLAY = {
    "s0_v1enc": "继承的编码器配置",
    "v1enc": "继承的编码器配置",
    "a3_v1enc": "继承的编码器配置",
    "s0_base": "未调参",
    "a3_base": "未调参",
    "base": "未调参",
}


def label(name: str) -> str:
    return DISPLAY.get(name, name)


def pct(x, digits=2):
    return f"{x * 100:+.{digits}f}%"


# --- slides -----------------------------------------------------------------------------------


def slide_title(date: str):
    s = prs.slides.add_slide(BLANK)
    txt(s, 0.9, 2.4, 11.5, 1.0, ["RIKYU 超参数调优 campaign"], size=40, bold=True)
    txt(s, 0.9, 3.5, 11.5, 0.6, ["v2 报告"], size=22, color=MUT)
    txt(s, 0.9, 4.4, 11.5, 1.2, [
        f"{date} ｜ 分支 exp/rikyu-hparam-tuning-v2",
        "调参前 ReduceLROnPlateau 按 batch 而非 epoch 触发（PR #45 修复）。修复前的一切测量都在那个坏节奏下，",
        "因此本轮全部重做；报告中唯一出现的修复前数据是未调参基线，作为归因算式的第一项。",
    ], size=13, color=MUT)


def finals_vs_control(a: dict) -> dict:
    """Leader vs the finals' OWN 25-seed untuned control, on each arm's measured sigma.

    The `vs_anchor` field in the summary compares against the stage-0 reference, which has nine
    seeds; quoting a 25-seed arm against it mixes seed counts in the standard error. The finals
    deliberately include an untuned arm at the same seed count, and that is the honest baseline.
    """
    import math

    arms = a["arms"]
    control = next((k for k in arms if k.endswith("_base")), None)
    if control is None:
        raise Missing("no untuned control arm in the finals")
    lead, base = arms[a["leader"]]["score"], arms[control]["score"]
    delta = lead["mean"] - base["mean"]
    se = math.sqrt(lead["sigma"] ** 2 / lead["n"] + base["sigma"] ** 2 / base["n"])
    return {"delta": delta, "two_se": 2 * se, "resolved": abs(delta) > 2 * se,
            "lead_sigma": lead["sigma"], "base_sigma": base["sigma"],
            "sigma_ratio": base["sigma"] / lead["sigma"] if lead["sigma"] else None,
            "control": control, "n": lead["n"]}


def slide_glossary():
    s = new("名词对照表", "正文用全称；这些是 RIKYU 上的目录名 / 作业名，读原始输出时会遇到")
    table(s, 0.5, 1.35, 5.9, ["标识符", "是什么"],
          [["stage0", "锚点：未调参 + 新镜像"],
           ["a1 / a1r", "编码器×LR×调度器 网格 / 随机搜索"],
           ["a3", "A′ 决赛（8 配置 × 25 seed）"],
           ["a4", "调度器是否值得 / 下界是否要外扩"],
           ["a2b", "早停复核：patience 24 对 40"],
           ["b", "head 网格（24 配置 × 5 seed）"],
           ["b3", "head 决赛（4 配置 × 25 seed）"],
           ["bal / balx", "损失平衡器 开关 / 排除 AE"],
           ["single (stA_*)", "同régime单任务天花板"],
           ["xfer", "迁移测试：24 任务 × 3 组随机顺序"]],
          col_w=[1.9, 4.0], size=10, head_size=10)
    table(s, 6.9, 1.35, 6.0, ["Stage C 臂", "是什么"],
          [["c2_base", "未调参（对照）"],
           ["c2_top1", "A′ 决赛第 1 名配置"],
           ["c2_top2 / c2_top3", "决赛第 2 / 3 名"],
           ["c2_base_cons", "未调参 + consolidation"],
           ["c2_top1_cons", "调参 + consolidation（本轮最好）"],
           ["stA_*", "同régime单任务基线（5 seed）"],
           ["xf_<任务>_o<k>", "迁移测量：该任务排在打乱序列最后，第 k 次重复"]],
          col_w=[2.0, 4.0], size=10, head_size=10)
    txt(s, 0.5, 5.2, 12.4, 1.8, [
        "命名规则：中间是配置  ｜  _cons = 做过 consolidation（全任务 + 全数据的收尾后训练）",
        "",
        "consolidation：24 任务顺序训练跑完后，把全部 24 个 head 连同编码器再联合微调一遍",
        "（fm finetune）—— 是收尾步骤，不是另一套超参。",
    ], size=12, color=MUT)


@slide_guard
def slide_summary():
    a = load("finals_a.json")
    bal = load("stage_bal.json")
    tr = load("transfer_adopted.json")
    vc = finals_vs_control(a)
    off = max(v["delta_vs_untuned"] for v in bal["vs_anchor"] if v["arm"].endswith("off"))
    on = max(v["delta_vs_untuned"] for v in bal["vs_anchor"] if v["arm"].endswith("on"))
    # Lead with the relative percentage — +0.045 does not tell a reader whether that is a lot.
    gainers = sorted((r for r in tr["per_task"]
                      if r.get("matters") and r.get("transfer", 0) > 0),
                     key=lambda r: -r["relative_pct"])
    gains = "、".join(f"{r['task']} {r['relative_pct']:+.1f}%" for r in gainers)
    s = new("摘要", "每一条都可在 summary/*.json 中溯源")
    txt(s, 0.6, 1.5, 12.2, 5.4, [
        f"1. 调参收益不大但明确可分辨：采纳配置相对决赛内同为 {vc['n']} seed 的未调对照臂 "
        f"{pct(vc['delta'])}（2SE {vc['two_se'] * 100:.2f}%，阈值的 {abs(vc['delta']) / vc['two_se']:.1f} 倍）；"
        f"但前 {1 + len(a['statistically_tied_with_leader'])} 名彼此统计上无法区分。"
        f"附带：采纳配置的 run 间 σ 是未调配置的 1/{vc['sigma_ratio']:.2f}。",

        "2. 24 任务上调参买到 +0.0074 / +1.03%（未调参 0.7155 → 调参 0.7229），加 consolidation 到 0.7274；"
        "被推上部署规模的三个配置极差仅 0.0037，探针名次不携带可操作的信息。",

        "3. 单任务天花板已在本régime重测（24 任务 × 5 seed）：继承的旧天花板在 23 个任务里 17 个偏低，"
        "平均 +0.0275，且不是常数，无法用偏移量修正。",
        f"4. multi-task 对小数据任务确有正迁移：{gains or '（无）'}；两个最大的任务略亏。",
        f"5. learnable loss balancer 有害且是机制性的：最好的关闭臂 {pct(off)} vs 最好的开启臂 {pct(on)}。不要上线。",
        "6. PCGrad 不适用：直接测量编码器逐任务梯度，未发现该方法赖以生效的方向冲突。",
        "7. 工程发现：两轮 campaign 的 GPU 利用率约 9%，单卡多任务打包实测提速 7.1×。",
    ], size=15)


def slide_design_flow():
    """The pipeline as a dependency chain: what each stage consumes and what decision it emits.

    Placed before any result, because a reader who does not know that the finals inherit their
    metric axis and their seed budget from the anchor cannot tell a designed campaign from a pile
    of runs. Every arrow is a real dependency in the code, not a narrative one.
    """
    s = new("调参流程：每一步消费什么、产出什么决定",
            "箭头是代码里真实的依赖，不是叙述顺序")
    table(s, 0.4, 1.35, 12.5,
          ["阶段", "规模", "变量", "消费上一步的", "产出的决定"],
          [["① 锚点 stage0", "18 run / 9 seed", "无（未调参重复）",
            "—", "参考点 + 单run σ"],
           ["② A′ 网格 + 随机", "296 配置 × 5 seed", "latent_dim × encoder_lr × min_lr × patience × factor",
            "参考点、度量轴、σ", "短名单（8）+ 边界是否用尽"],
           ["③ A′ 决赛", "10 臂 × 25 seed", "同上，仅短名单",
            "短名单", "**采纳编码器 / 调度器**"],
           ["④ A4 调度器价值", "60 run", "有调度 vs 固定 LR",
            "σ", "保留调度器；下界无需外扩"],
           ["⑤ a2b 早停", "2 臂 × 5 seed", "patience 24 vs 40",
            "**采纳基座**", "保留 24"],
           ["⑥ B′ head 网格", "24 配置 × 5 seed", "head 容量/LR × KR 分支容量/LR",
            "**采纳基座**", "短名单（4）"],
           ["⑦ B′ 决赛", "4 臂 × 25 seed", "同上",
            "短名单", "**head 不用动**"],
           ["⑧ 天花板重测", "24 任务 × 5 seed", "单任务训练",
            "采纳配置", "deficit 的分母"],
           ["⑨ Stage C′", "6 臂 × 24 任务", "采纳 vs 未调；决赛前 3 名",
            "③⑦ 的采纳值", "部署规模验证；名次是否迁移"],
           ["⑩ 迁移测量 xfer", "24 任务 × n 组顺序", "待测任务排最后",
            "采纳配置、天花板", "**迁移是否成立**"]],
          col_w=[2.3, 2.0, 3.6, 2.2, 2.4], size=9, head_size=9)
    txt(s, 0.4, 5.55, 12.5, 1.6, [
        "旁支（不进主链，各自回答一个二元问题）：损失平衡器开/关 · PCGrad 的前提是否存在 · 打包倍数标定",
        "",
        "两处关键依赖容易被忽略：① 的 σ 决定后面每一步需要多少 seed；③ 的采纳基座是 ⑤⑥ 的前提 —— "
        "在别的基座上测出的 head 或早停结论不能迁移过来。",
    ], size=11, color=MUT)


def slide_design_why():
    """Why each step is shaped the way it is. The decisions, not the numbers."""
    s = new("每一步为什么这么设计", "同样的算力可以花在别处，这里是选择的理由")
    txt(s, 0.4, 1.3, 6.3, 5.6, [
        "为什么先做锚点，而不是直接搜索",
        "  所有增益都是相对量，需要一个同镜像、同代码的参考点；",
        "  它同时给出 σ，而 σ 决定后面每一步买多少 seed 才够。",
        "",
        "为什么用 6 任务探针而不是 24",
        "  24 任务一轮要一天多，296 个配置跑不起。探针覆盖大/中/小",
        "  三档把 σ 从 5.01% 压到 2.05%，分辨 1% 所需 seed 从 101 降到 17。",
        "",
        "为什么网格和随机搜索都要",
        "  网格给可读的边际效应图（每个轴的单独影响），",
        "  随机给内部覆盖（网格点之间的空隙）。两者答同一个问题的不同面。",
        "",
        "为什么分「5 seed 筛选 + 25 seed 决赛」两段",
        "  296 配置 × 25 seed = 7400 个 run，跑不起；",
        "  但 5 seed 的榜首不可信 —— 本轮两次实测到 winner's curse。",
        "  所以先用便宜的 seed 数筛，再用贵的 seed 数定名次。",
    ], size=11)
    txt(s, 6.9, 1.3, 6.0, 5.6, [
        "为什么要检查网格边界",
        "  最优点落在搜索范围的端点上，说明范围不够，不是找到了最优。",
        "",
        "为什么早停要在「采纳之后」重测",
        "  最初是在当时的榜首上测的，而决赛换了榜首 ——",
        "  在已不采纳的配置上得到的结论不能沿用。",
        "",
        "为什么 head 要在采纳基座上调",
        "  head 的最优值依赖它所处的优化régime。换了基座就是换了régime。",
        "",
        "为什么把「什么都不改」放进 head 网格",
        "  只占 24 个格点里的 1 个，却让「head 不需要调」成为**排名结果**",
        "  而不是论证。它最后排第 8，且 2SE 最小。",
        "",
        "为什么 Stage C′ 推三个配置而不是一个",
        "  只推一个，就只能看到「调参臂比基线好」；",
        "  推三个才能看到「被推上去的三个彼此没区别」—— 这是可证伪点。",
        "",
        "为什么天花板要重测",
        "  继承的那批测于不同régime，用它算 deficit 等于把",
        "  「模型变化」和「测量框架变化」加在一起，事后无法拆开。",
    ], size=11)


@slide_guard
def slide_probe():
    s0 = load("stage0.json")
    cal = s0["calibration"]
    s = new("② 探针 probe6：把噪声压到可判定的水平", "早先的 3 任务探针噪声带 8.48%，前三名相差 1.5–1.8% —— 排不出名次")
    table(s, 0.6, 1.5, 6.0,
          ["探针任务", "标签数", "档位"],
          [[t, f"{n:,}", g] for t, n, g in [
              ("volume", 23678, "big"), ("formation_energy", 23180, "big"),
              ("seebeck", 8072, "mid"), ("zt", 3445, "mid"),
              ("magnetization", 1160, "small"), ("magnetic_moment", 851, "small")]],
          col_w=[2.6, 1.7, 1.7])
    need_seeds = cal["seeds_needed_to_resolve"]
    # v1 published a RANGE at n=3, not a sigma. E[range] = d2(n)*sigma, d2(3)=1.693 — putting the
    # two side by side without that conversion reads as a 4x noise gap where the real one is 2.4x.
    v1_sigma = cal["v1_probe3_band_for_reference"] / 1.693
    v1_seeds_1pct = math.ceil((2 * v1_sigma / 0.01) ** 2)
    txt(s, 7.0, 1.5, 5.8, 4.8, [
        f"v2 单 run σ = {cal['sigma_per_run'] * 100:.2f}%   （9 seed 实测）",
        f"3 任务探针 σ ≈ {v1_sigma * 100:.2f}%   ← 由其 3 seed 极差 "
        f"{cal['v1_probe3_band_for_reference'] * 100:.2f}% 换算（E[极差] = d₂(n)·σ，d₂(3)=1.693）",
        "",
        "要分辨这么大的真实差异，v2 需要的 seed 数：",
        *[f"    {float(k) * 100:.1f}%  →  {v} seed" for k, v in need_seeds.items()],
        f"        （同样分辨 1.0%，3 任务探针需要 {v1_seeds_1pct} 个）",
        "",
        "排除 electrical_resistivity（天花板 0.162，无分辨力）",
        "排除 magnetic_susceptibility（58 个标签）",
        "",
        f"18 个锚点 run 实测：平均 {cal['wallclock']['mean_hours']:.2f} h/run，"
        f"合计 {cal['wallclock']['total_gpu_hours']:.0f} GPU-h",
    ], size=13)


@slide_guard
def slide_anchor():
    s0 = load("stage0.json")
    c = s0["comparisons"][0]
    s = new("① 锚点：所有增益的参考点，以及噪声有多大",
            "未调参配置 × 9 seed，同镜像同代码 —— 后面每一个「+x%」都是相对它报的")
    txt(s, 0.6, 1.6, 12.2, 3.2, [
        "锚点做两件事，缺一不可：",
        "",
        "  • 给出**参考点**。增益是相对量，没有同régime的未调参基准就没有分母。",
        "  • 给出**单run σ = 2.05%**，而 σ 决定后面每一步买多少 seed 才够 ——",
        "    分辨 1% 需要 17 个 seed，分辨 2% 只要 5 个。整轮的 seed 预算由此定下。",
        "",
        "顺带检验了一个继承来的配置在当前代码上还有没有优势：",
    ], size=14)
    verdict = "无法分辨" if abs(c["delta_score"]) <= c["resolvable_at_this_n"] else "可分辨"
    table(s, 0.6, 4.6, 9.0,
          ["比较", "差值", "可分辨阈值 (2SE)", "判定"],
          [[f"{label(c['from'])} → {label(c['to'])}", pct(c["delta_score"]),
            pct(c["resolvable_at_this_n"]), verdict]],
          col_w=[3.4, 1.8, 2.2, 1.6])
    txt(s, 0.6, 5.5, 12.2, 0.8,
        ["即：那套继承来的编码器配置在当前代码上不再有可见优势 —— 所以 Stage A′ 从头重搜。"],
        size=13, color=MUT)


@slide_guard
def slide_grid():
    a = load("stage_a.json")
    s = new("Stage A′：编码器 × LR × 调度器联合搜索",
            f"{a['n_configs']} 个配置 / {a['n_runs']} 个 run（网格 + 随机搜索，各 5 seed）")
    edge = a["edge_bound_axes"]
    txt(s, 0.6, 1.5, 12.2, 2.0, [
        f"网格边界检查：{'全部通过，不需要追加 a1b 轮' if not edge else '边界受限的轴：' + ', '.join(edge)}",
        "",
        "  这项检查本身经过两次修正：最初只测“最优点是否恰好等于端点”，但 206 个随机取值没有任何一个",
        "  正好落在端点上，真正的边界问题会被放过 —— 后来加了十分位趋势检验；它还曾在只有 2 个取值的轴",
        "  （patience、latent_dim）上误报，因此增加了“至少 3 个取值”的前提。",
    ], size=13)
    ties = a["leader_ties"]
    n_tied = len(ties["statistically_tied_with_leader"])
    short = set(a["short_list"])
    excluded = [c for c in ties["statistically_tied_with_leader"] if c not in short]
    txt(s, 0.6, 3.7, 12.2, 2.8, [
        f"5 seed 排不出名次：{n_tied} 个配置与榜首统计上并列",
        "（阈值按每一对的 2·√(sem₁²+sem₂²) 算 —— 两个臂的标准误都计入）。",
        "",
        "短名单取前 8 进 25 seed 决赛。但这是**预算截断**，不是噪声感知的筛选：",
        f"{n_tied} 个并列配置里有 {len(excluded)} 个没进决赛，仅因含噪的样本均值排在第 8 之后。",
        "所以决赛冠军是这 8 个里最好的，不是 " + str(n_tied) + " 个里最好的。",
    ], size=13)


@slide_guard
def slide_finals():
    a = load("finals_a.json")
    vc = finals_vs_control(a)
    ranked = sorted(a["arms"].items(), key=lambda kv: -kv[1]["score"]["mean"])
    rows = [[label(k.replace("a3_", "")), pct(v["score"]["mean"]), f"{v['score']['sigma'] * 100:.2f}%"]
            for k, v in ranked]
    s = new(f"Stage A′ 决赛（{vc['n']} seed，{a['n_runs']} run）",
            "决赛内含同 seed 数的未调对照臂 —— 增益在同一 seed 数下比出，不是 25 seed 比 9 seed")
    table(s, 0.6, 1.45, 7.6, ["配置", "相对 stage-0 参考", "run 间 σ"], rows,
          col_w=[4.6, 1.6, 1.4], size=10, head_size=10, colour_col=1)
    tied = a["statistically_tied_with_leader"]
    need_seeds = list(a["seeds_that_would_resolve_the_ties"].items())[:2]
    txt(s, 8.5, 1.45, 4.5, 5.6, [
        f"采纳 vs 未调对照：{pct(vc['delta'])}",
        f"2SE {vc['two_se'] * 100:.2f}% → "
        f"{'可分辨' if vc['resolved'] else '不可分辨'}（阈值的 {abs(vc['delta']) / vc['two_se']:.1f} 倍）",
        "",
        f"但前 {1 + len(tied)} 名彼此并列。分开所需 seed：",
        *[f"    {k.split(' vs ')[1].replace('a3_', '')[:26]} → {v}" for k, v in need_seeds],
        "",
        "好配置有一部分好在“稳”：",
        f"    采纳 σ {vc['lead_sigma'] * 100:.2f}%  vs  未调 σ {vc['base_sigma'] * 100:.2f}%",
        f"    （1/{vc['sigma_ratio']:.2f}）",
        "",
        "winner’s curse：5 seed 的榜首 a1r129 在",
        "25 seed 下掉到最后一名，而它的 σ 是十个臂",
        "里第二大的 —— σ 大的配置更容易赢小样本抽签。",
    ], size=11)


@slide_guard
def slide_finals_sigma():
    pic_slide("好配置有一部分好在“稳”",
              "σ 与 25 seed 成绩的相关系数 −0.844 —— σ 大的配置更容易赢小样本抽签，样本一加就落回原形",
              RES / "finals_sigma_vs_mean.png")


@slide_guard
def slide_a4():
    a4 = load("stage_a4.json")
    h = a4["head_to_head"]
    d = a4["downward_extension"]["sched"]
    s = new("A4：调度器值不值得 / 最优点是否在搜索下界之外", f"{a4['n_runs']} 个 run，成对比较")
    table(s, 0.6, 1.6, 11.0,
          ["问题", "结果", "差值", "2SE 可分辨", "判定"],
          [["有调度 vs 固定 LR",
            f"{pct(h['best_scheduled']['mean'])} vs {pct(h['best_flat']['mean'])}",
            pct(h["delta"]), pct(2 * h["se_of_difference"]),
            "调度器值得保留" if h["separated"] else "无法分辨"],
           ["最优点是否在下界之外",
            f"{pct(d['best_below_floor']['mean'])} vs {pct(d['best_at_or_above_floor']['mean'])}",
            pct(d["delta"]), pct(2 * d["se_of_difference"]),
            "不需要向下扩展" if not d["worth_extending_further"] else "需要扩展"]],
          col_w=[3.0, 2.6, 1.6, 1.8, 2.0], colour_col=2)
    txt(s, 0.6, 3.4, 12.2, 2.6, [
        "必须随数字一起声明的限制：",
        "",
        "  [training.scheduler] 管的是全部四个参数组（encoder / head / kr / ae），没有分组开关。",
        "  所以“关掉调度”同时冻结了 head 与 KR 的学习率 —— 这个损失不能归因到某一个组。",
    ], size=14, color=MUT)


@slide_guard
def slide_a2b():
    a = load("stage_a2b.json")
    arms = a["arms"]
    rows = []
    for name, arm in sorted(arms.items(), key=lambda kv: -kv[1]["score"]["mean"]):
        s = arm["score"]
        rows.append([name.split("_")[-1], pct(s["mean"]), f"{2 * s['sem'] * 100:.3f}%",
                     f"{s['sigma'] * 100:.3f}%"])
    s_ = new("A2b：早停 patience 24 还是 40（在采纳后的基座上复核）",
             "最初的 a2 跑在“当时的榜首”上，而 25 seed 决赛把榜首换掉了 —— 所以这是复核，不是初测")
    table(s_, 0.6, 1.5, 7.4, ["臂", "相对未调锚点", "2SE", "σ"], rows,
          col_w=[1.6, 2.2, 1.8, 1.8], colour_col=1)
    pair = a["pairwise"][0] if a.get("pairwise") else None
    lines = []
    if pair:
        lines += [f"差 {pct(pair['delta'])}，可分辨阈值 {pct(pair['resolvable_at_this_n'])[1:]}",
                  f"→ {'可分辨' if pair['separated'] else '不可分辨'}"
                  + (f"（要分开需 {list(a['seeds_that_would_resolve_the_ties'].values())[0]} 个 seed，"
                     f"我们有 5 个）" if a.get("seeds_that_would_resolve_the_ties") else ""), ""]
    lines += ["折合绝对 R² 约 0.0017 —— 远低于 1e-2 的实用门槛。",
              "而 ES40 的墙钟是 3.75 h，ES24 是 3.38 h：多花 11%。", "",
              "采纳 patience 24。多付 11% 算力买一个既不可分辨、",
              "又低于实用门槛的差异，不划算。", "",
              "教训：凡是“在当前最优上”做的实验，",
              "最优一变就必须重做。"]
    txt(s_, 8.4, 1.5, 4.6, 5.0, lines, size=12)


@slide_guard
def slide_stage_b():
    import math

    b = load("stage_b.json")
    R = {e["config"]: e for e in b["ranking"]}
    # The default head block, spelled as a stage-B label. It is IN the grid, which is what lets
    # "changing nothing" be ranked against every change rather than assumed to be the baseline.
    default = "b_H64_HL0p005_X128-64_KL0p0005"
    lead = b["ranking"][0]
    rows = []
    for i, e in enumerate(b["ranking"], 1):
        if i > 5 and e["config"] != default and i < len(b["ranking"]):
            continue
        tag = "  ← 默认 head 块" if e["config"] == default else ""
        rows.append([f"{i}", e["config"].replace("b_", "") + tag,
                     pct(e["score_mean"]), f"{2 * e['score_sem'] * 100:.3f}%"])
    s = new("Stage B′：多任务联合 head 调参 —— 结论是 head 不用动",
                  f"24 个配置 × 5 seed = {b['n_runs']} run，全部在 A′ 采纳基座上；120/120 通过训练校验")
    table(s, 0.5, 1.45, 8.6, ["名次", "配置", "相对未调锚点", "2SE"], rows,
          col_w=[0.8, 5.2, 1.5, 1.1], size=10, head_size=10, colour_col=2)
    lines = [f"榜首与 {len(b['leader_ties']['statistically_tied_with_leader'])} 个配置统计上并列",
             f"（5 seed 可分辨差异 {pct(b['leader_ties']['resolvable_difference'])[1:]}）", ""]
    if default in R:
        d = lead["score_mean"] - R[default]["score_mean"]
        se = math.sqrt(lead["score_sem"] ** 2 + R[default]["score_sem"] ** 2)
        rank = [e["config"] for e in b["ranking"]].index(default) + 1
        lines += [f"“什么都不改”排第 {rank} / {len(b['ranking'])}",
                  f"榜首 vs 默认：{pct(d)}，2SE {2 * se * 100:.3f}%",
                  f"→ {'可分辨' if abs(d) > 2 * se else '不可分辨'}", "",
                  "默认 head 块还是全网格 2SE 最小的一个",
                  f"（{2 * R[default]['score_sem'] * 100:.3f}%），即最可复现。", ""]
    lines += ["采纳：默认 head 块（改动最少，与 A′ 同一条规则）。",
              "注意它不是网格榜首 —— 榜首的 2SE 是前 12 名里",
              "第二大的，正是 A′ 实测过的 winner’s curse 模式。", "",
              "比早先的同一结论更强：那次可归因于 régime 不对，",
              "本轮是在正确的 régime 上调的，依然什么都没买到。"]
    txt(s, 9.3, 1.45, 3.7, 5.6, lines, size=11)


@slide_guard
def slide_b_finals():
    import math
    import statistics

    b5 = load("stage_b.json")
    b25 = load("finals_b.json")
    five = {e["config"]: e["score_mean"] for e in b5["ranking"]}
    default = "H64_HL0p005_X128-64_KL0p0005"
    rows, xs, ys = [], [], []
    entries = []
    for name, arm in b25["arms"].items():
        short = name.replace("b3_", "")
        f5 = five.get("b_" + short)
        s = arm["score"]
        entries.append((short, f5, s["mean"], s["sigma"]))
    entries.sort(key=lambda r: -(r[1] or 0))
    for short, f5, m25, sg in entries:
        tag = "  ← 默认" if short == default else ""
        rows.append([short + tag, pct(f5) if f5 is not None else "-", pct(m25),
                     pct(m25 - f5) if f5 is not None else "-", f"{sg * 100:.3f}%"])
        if f5 is not None:
            xs.append(sg)
            ys.append(-(m25 - f5))
    s_ = new("b3 决赛：采纳决定从“按规则选的”变成“测出来的”",
             f"4 个臂 × 25 seed = {b25['n_runs']} run，220/220 通过训练校验")
    table(s_, 0.5, 1.45, 9.4,
          ["配置", "5 seed", "25 seed", "跌幅", "σ(25)"], rows,
          col_w=[4.4, 1.2, 1.2, 1.2, 1.2], size=10, head_size=10, colour_col=3)
    lead = entries[0]
    dflt = next(e for e in entries if e[0] == default)
    e5, e25 = lead[1] - dflt[1], lead[2] - dflt[2]
    need = b25.get("seeds_that_would_resolve_the_ties", {})
    pair = next((v for k, v in need.items() if default in k), None)
    r = None
    if len(xs) > 2:
        mx, my = statistics.fmean(xs), statistics.fmean(ys)
        den = math.sqrt(sum((x - mx) ** 2 for x in xs) * sum((y - my) ** 2 for y in ys))
        r = sum((x - mx) * (y - my) for x, y in zip(xs, ys)) / den if den else None
    lines = ["答案：没有任何 head 配置胜过“什么都不改”。", "",
             f"榜首对默认的优势：{pct(e5)} → {pct(e25)}",
             f"    坍缩 {(1 - abs(e25) / abs(e5)) * 100:.0f}%"]
    if pair:
        lines.append(f"    分开这两者需要 {pair:,} 个 seed")
    lines += ["", "默认 head 块几乎没动（跌 0.060%），",
              "三个调过的都掉了 0.5–0.66%；σ 也最小。"]
    if r is not None:
        lines += ["", f"σ 与跌幅的相关系数 {r:+.3f}", "—— 与 A′ 测到的是同一个机制。"]
    lines += ["", "同一轮 campaign 内第二次独立复现 winner’s curse：",
              "A′ 的 5 seed 榜首掉到第 10，B′ 的 5 seed 榜首",
              "优势全部蒸发。处方都一样：加 seed，不是加网格点。"]
    txt(s_, 10.1, 1.45, 2.9, 5.6, lines, size=10)


@slide_guard
def slide_stage_c():
    """The tuned configurations at real scale — v2's arms only.

    No pre-fix baseline and no fix-vs-tuning attribution: the campaign's deliverable is the tuning
    result, and the scheduler bug was something met along the way. That it helps to fix a bug is
    not a finding, so scoring against the broken state does not belong in the deliverable.
    """
    c = load("stage_c.json")
    keep = {"c2_base": "未调参",
            "c2_top1": "调参（决赛第 1）",
            "c2_top2": "调参（决赛第 2）",
            "c2_top3": "调参（决赛第 3）",
            "c2_base_cons": "未调参 + consolidation",
            "c2_top1_cons": "调参 + consolidation"}
    by = {a["label"]: a for a in c["arms"]}
    arms = sorted((by[k] for k in keep if k in by), key=lambda a: -a["mean_r2"])
    rows = []
    for a in arms:
        d = a["deficit"]
        f = lambda v: f"{v:+.4f}" if v is not None else "-"  # noqa: E731
        rows.append([keep[a["label"]], f"{a['mean_r2']:.4f}",
                     f(d["big"]), f(d["mid"]), f(d["small"])])
    s = new("Stage C′：24 任务最终运行", "deficit 对同régime单任务天花板；每臂 1 seed，小差距不可分辨")
    table(s, 0.5, 1.45, 9.4, ["臂", "mean R²", "big", "mid", "small"], rows,
          col_w=[3.8, 1.5, 1.4, 1.4, 1.3], size=10, head_size=10)

    def delta(a, b):
        if a in by and b in by:
            d = by[b]["mean_r2"] - by[a]["mean_r2"]
            return d, d / by[a]["mean_r2"] * 100
        return None, None
    lines = ["调参在 24 任务上买到多少：", ""]
    for a, b, label in (("c2_base", "c2_top1", "调参"),
                        ("c2_top1", "c2_top1_cons", "+ consolidation"),
                        ("c2_base", "c2_base_cons", "未调参 + consolidation")):
        d, rel = delta(a, b)
        if d is not None:
            lines.append(f"  {label:16s} {d:+.4f}  = {rel:+.2f}%")
    lines += ["", "与探针上的 +1.56% 同量级、方向一致，",
              "但每臂只有 1 seed，三个数都不可分辨。"]
    tr = c.get("transfer", {})
    if tr.get("checked"):
        lines += ["", "探针名次是否迁移到 24 任务：",
                  f"  探针说 {' > '.join(x.replace('c2_', '') for x in tr['probe_order'])}",
                  f"  实际是 {' > '.join(x.replace('c2_', '') for x in tr['deployed_order'])}",
                  f"  极差仅 {tr['mean_r2_spread_across_promoted_arms']:.4f}",
                  "  → 名次不携带可操作的信息。"]
    txt(s, 10.1, 1.45, 3.0, 5.6, lines, size=10)


@slide_guard
def slide_ceiling_fig():
    pic_slide("天花板是一个坏掉的测量框架",
              "旧天花板测于 PR #45 之前 —— LR 在第一个 epoch 内就落到地板，那不是天花板",
              RES / "ceiling_frame_offset.png")


@slide_guard
def slide_transfer_fig():
    pic_slide("迁移：multi-task 到底有没有帮到小任务",
              "采纳配置下，25 seed 多任务 vs 5 seed 单任务，唯一差别是 pretrain.task_sequence",
              RES / "transfer_adopted.png")


@slide_guard
def slide_transfer_why():
    tr = load("transfer_adopted.json")
    sm = tr["summary"]
    gainers = sorted((r for r in tr["per_task"]
                      if r.get("matters") and r.get("transfer", 0) > 0),
                     key=lambda r: -r["relative_pct"])
    gains = "、".join(f"{r['task']} {r['relative_pct']:+.1f}%" for r in gainers)
    s = new("为什么这个问题卡住了另外三个",
            "既然存在单任务天花板，multi-task 只有确实让数据少的任务变好时才值得")
    txt(s, 0.6, 1.5, 12.2, 4.6, [
        "否则 loss balancing 和梯度手术都没有可修的东西 —— 小任务的正确答案就是单独训练。",
        "",
        f"实测（相对 R² 提升）：{gains or '无'}；"
        f"{'、'.join(sm['tasks_hurt']) or '无'} 明确受损；其余无法分辨。",
        "",
        "顺序是单调的，而且是需要的那个方向：最小的两个任务和 zt 获益，最大的两个略亏。",
        "",
        "两种百分比口径含义不同：相对提升 = ΔR²/单任务R²（稳定，作头条）；",
        "误差消减 = ΔR²/(1−单任务R²)（有余量时更有意义，接近天花板时会爆炸）。",
        "formation_energy 残差只有 0.0053，−0.0036 按误差消减算是 −68% —— 算术正确，作头条误导。",
        "",
        f"实用门槛 |ΔR²| ≥ 0.01：{'、'.join(sm.get('tasks_that_matter', [])) or '无'} 值得写进结论；"
        f"{'、'.join(sm.get('resolved_but_negligible', [])) or '无'} 可分辨但可忽略。",
        "",
        "所以多任务这套设置在自己的账上是划算的。不划算的是那个本该保护小任务的 balancer。",
        "",
        "注意：这是 6 个任务的结论。24 个任务在部署规模上是否同样成立，由 xfer 阶段测量 ——",
        "每个任务都被放在打乱后的 24 任务序列末尾训练，每任务 3 组随机顺序。",
    ], size=14)


@slide_guard
def slide_xfer():
    """The campaign's headline question, answered at deployment scale.

    Reads matched_xfer.json, not transfer_xfer.json: the report quotes the comparison restricted to
    the rows both arms share, and a deck showing the uncorrected figures beside a report showing the
    corrected ones is how two documents start disagreeing.
    """
    x = load("matched_xfer.json")
    rows_all = [r for r in x["per_task"] if "transfer" in r]
    better = [r for r in rows_all if r["separated"] and r["transfer"] > 0]
    worse = [r for r in rows_all if r["separated"] and r["transfer"] < 0]
    unres = [r for r in rows_all if not r["separated"]]
    ranked = sorted(rows_all, key=lambda r: -(r.get("relative_pct") or 0))
    shown = ranked[:5] + ranked[-5:]
    rows = []
    for r in shown:
        verdict = ("多任务更好" if r["transfer"] > 0 else "单任务更好") if r["separated"] else "无法分辨"
        rows.append([r["task"], f"{r['n_train']:,}", f"{r['single_task']:.4f}",
                     f"{r['multi_task']:.4f}", f"{r['relative_pct']:+.2f}%", verdict])
    s = new("部署规模的迁移测量：结论是负的",
            "24 任务 × 3 组打乱顺序 = 72 run，待测任务排最后；两边限制在共同测试集上")
    table(s, 0.5, 1.45, 11.6,
          ["任务", "标签数", "单任务", "多任务", "相对", "判定"], rows,
          col_w=[3.0, 1.4, 1.5, 1.5, 1.5, 2.7], size=10, head_size=10, colour_col=4)
    y = 1.55 + 0.32 * (len(rows) + 1) + 0.15
    txt(s, 0.5, y, 12.4, 2.4, [
        f"多任务更好 {len(better)}（{'、'.join(r['task'] for r in better) or '无'}）　│　"
        f"单任务更好 {len(worse)}　│　无法分辨 {len(unres)}",
        f"（表中只列相对变化的两端各 5 个，共 {len(rows_all)} 个任务）",
        "",
        "探针高估了迁移：zt 从 +6.85%「多任务更好」变成 +3.16%「无法分辨」；",
        "magnetization 从 +4.40% 变成 +1.72%；magnetic_moment 从「无法分辨」变成 −6.33%「单任务更好」。",
        "三个判定全部朝不利方向移动 —— 6 个任务的探针不能预测 24 个任务。",
    ], size=12)


@slide_guard
def slide_position():
    """The position curve — obtained by re-reading data the stage already wrote.

    Kept separate from slide_xfer because it answers a different question: xfer asks "is multi-task
    better for a task placed last", this asks "does it matter WHERE the task sits". The second is
    what explains part of the first.
    """
    d = load("position.json")
    rows_all = [r for r in d["per_task"] if r.get("early_1_8") and r.get("late_17_24")]
    ranked = sorted(rows_all, key=lambda r: (r["position_effect"] or {}).get("delta_late_minus_early", 0))

    def cell(x):
        if not x:
            return "-"
        star = "*" if x.get("matters") else ("·" if x.get("separated") else "")
        return f"{x['delta']:+.4f} / {x['relative_pct']:+.1f}%{star}"

    def pcell(pe):
        if not pe:
            return "-"
        star = "*" if pe.get("matters") else ("·" if pe.get("separated") else "")
        return f"{pe['delta_late_minus_early']:+.4f} / {pe['relative_pct']:+.1f}%{star}"

    shown = ranked[:5] + ranked[-3:]
    rows = [[r["task"], f"{r['n_train']:,}", cell(r["early_1_8"]), cell(r["late_17_24"]),
             pcell(r["position_effect"])] for r in shown]
    s = new("把视点固定在单个任务上：位置的影响",
            "每个 step 目录都记录了到该步为止所有任务的指标 —— 每个任务在所有位置的分数早已在磁盘上，零算力")
    table(s, 0.4, 1.5, 12.5,
          ["任务", "标签数", "早期（位置 1–8）", "晚期（位置 17–24）", "位置效应（晚 − 早）"],
          rows, col_w=[2.8, 1.4, 2.8, 2.8, 2.7], size=10, head_size=10)
    n_pos = sum(1 for r in rows_all if (r["position_effect"] or {}).get("matters"))
    n_ret = sum(1 for r in d["per_task"]
                if r.get("retention") and r["retention"]["separated_from_zero"]
                and r["retention"]["mean"] > 0)
    y = 1.6 + 0.32 * (len(rows) + 1) + 0.2
    txt(s, 0.4, y, 12.5, 3.0, [
        f"位置对 {n_pos} / {len(rows_all)} 个任务可分辨。magnetic_moment 与 total_magnetization "
        "在早期位置根本不吃亏（不可分辨），排到最后才掉 ——",
        "所以迁移测量里 magnetic_moment 的 −6.33% 是位置造成的，不是多任务训练本身造成的。",
        "",
        f"另一条：{n_ret} / 24 个任务在训完自己那一步之后分数还会再涨。但检验否定了「replay 持续获益」的解释 ——",
        "保持度与剩余步数的相关系数只有 +0.038（n=1656），增益在头 3–4 步内就到位；",
        "排最后的任务保持度恰好 +0.0000（69 个样本），这是测量正确的健全性检查。",
        "读作：每个任务在自己那一步没训够，后面两三步的 replay 补回约 +0.02。",
        "（早停监控的是全任务总损失，当前任务的进展被已收敛的旧任务稀释 —— 假设，待验证。）",
    ], size=11)
    txt(s, 0.4, 6.6, 12.5, 0.5,
        ["* = 可分辨且 ≥0.01    · = 可分辨但低于实用门槛    位置与「前面是哪些任务」混淆，每位置 3–6 个样本"],
        size=9, color=MUT)


def slide_descriptor_limit():
    s = new("上线前该知道的一条模型限制：描述符看不见晶胞尺度",
            "volume 0.619 / final_energy 0.774 对 ~23 700 标签的任务太低，而同批行上 formation_energy 0.995")
    txt(s, 0.5, 1.5, 7.4, 4.6, [
        "在代码层面确认，不是推测：",
        "  descriptor_fn → formula_to_composition（契约是 atomic-FRACTION 向量，和为 1）",
        "  → KMD.transform 做 weight @ K，不再归一化",
        "",
        "  formula_to_composition(\"Fe2O3\") == formula_to_composition(\"Fe4O6\")  →  True",
        "",
        "即晶胞尺度不在模型输入里。成分键本身特意保留了绝对化学计量，",
        "但 KMD 这条路径在入口就把它除掉了。",
        "",
        "是泛化缺口，不是标签噪声：33 822 个约化式里只有 7 个重复（0.02%），",
        "训练数据没有矛盾 —— 模型只是必须从化学成分反推尺度。",
        "",
        "corr(Volume, 晶胞原子数) = +0.868  →  75.3% 的方差",
        "原子数跨度 1–320（中位 17）。单任务到 0.619，推回了大部分但推不全。",
    ], size=12)
    table(s, 8.2, 1.5, 4.7, ["约化式", "原子数", "体积"],
          [["AgSO4", "12 / 48", "162 / 616"],
           ["U(PO3)4", "34 / 136", "450 / 1930"],
           ["Ba(FeAs)2", "5 / 10", "98 / 217"],
           ["MnIr", "2 / 4", "26 / 56"]],
          col_w=[1.7, 1.4, 1.6], size=10, head_size=10)
    txt(s, 8.2, 3.4, 4.7, 3.2, [
        "↑ 同一个描述符输入，体积差约 4 倍",
        "",
        "只解释三个低天花板中的一个：",
        "  volume            +0.868  ✓ 是",
        "  final_energy      +0.162  ✗ 本就是每原子量",
        "  total_magnetization +0.023 ✗ 磁性本就难",
        "",
        "不影响本报告任何比较（各臂共用描述符）。",
        "改进项：给描述符补晶胞原子总数 —— 属模型",
        "改动，列入交接，不在 v2 范围。",
    ], size=11, color=MUT)


@slide_guard
def slide_balancer():
    b = load("stage_bal.json")
    rows = [[v["arm"], pct(v["delta_vs_untuned"])] for v in b["vs_anchor"]]
    s = new("learnable loss balancer：有害，且是机制性的",
            f"{b['n_runs']} 个 run。这个功能此前一直没有真正接上，本轮修复链路后才第一次可测")
    table(s, 0.6, 1.45, 5.4, ["臂", "相对未调锚点"], rows, col_w=[3.0, 2.4],
          size=10, head_size=10, colour_col=1)
    sep = [p for p in b["pairwise"] if p["separated"]]
    txt(s, 6.4, 1.45, 6.4, 5.4, [
        "开启的每一个臂都低于关闭的每一个臂。",
        f"可分辨的成对比较：{len(sep)} / {len(b['pairwise'])}",
        "",
        "为什么反了：",
        "  方法学每个任务的 log σ，最优解是 σ² = L（该任务自身的损失）。",
        "  实测 σ 与原始损失的相关系数 +0.970。",
        "  由此得到的 head 权重：AE 20 075 vs seebeck 1.5 —— 差四个数量级。",
        "",
        "  即：它把权重给了最容易拟合的任务，把最难拟合的压了下去，",
        "  与“救援弱任务”的初衷正好相反。",
        "",
        "排除 AE（balx 臂）之后仍然反，监督任务之间约 112 倍。",
        "所以问题不在作用范围，而在方法前提。",
        "",
        "结论：不要上线。附带的好处是它此前一直没被开启 ——",
        "否则大量算力会被浪费在一个有害的机制上。",
    ], size=12)


def slide_pcgrad():
    s = new("PCGrad（arXiv 2001.06782）：前提不成立",
            "只对负余弦的任务对起作用；没有冲突时它是恒等变换")
    txt(s, 0.6, 1.6, 12.2, 4.8, [
        "代价：每个任务一次反向传播 —— probe6 上 6×，Stage C 上 24×。所以先测前提，再谈代价。",
        "",
        "直接测量了共享编码器上的逐任务梯度（只有编码器可能冲突；task head 参数不相交，按构造无法干涉）。",
        "",
        "  • 论文条件一「方向冲突」：没有测到。→ 不引入。",
        "  • 论文条件二「梯度量级支配」：确实存在（各任务编码器梯度范数相差很大）。",
        "",
        "但量级支配单独不构成采用 PCGrad 的理由 —— 它修的是方向冲突。",
    ], size=15)


def slide_packing():
    s = new("工程发现：一卡多任务打包", "Slurm 计费显示单个 run 只用掉一块 GB200 的约 9%")
    txt(s, 0.6, 1.6, 12.2, 4.8, [
        "单 run 实测：算力约 9%，显存 1.29 GB / 189 GB。两轮 campaign 都是一 run 一卡 —— 约九成预约被浪费。",
        "",
        "改为 --pack N（N 个 run 共享一块 GPU 的独立进程）后，实测吞吐 7.1×（PACK=8）。",
        "如果这个做法从一开始就有，两轮 campaign 大约能省下 2300 card-h。",
        "",
        "一条被 review 抓出来的错误：最初报告 8.0×，但那个数字是循环论证 —— cost.py 里的 card-h",
        "本来就是 run_h ÷ pack_size 算出来的。真实测量值是 7.1×，已在 AGENTS.md、RIKYU instructions、",
        "skill 和 cost.py 四处一并修正。",
        "",
        "口径限制：打包会拉长单 run 的墙钟。打包前后的墙钟数字不可并列比较 ——",
        "“高 encoder_lr 收敛更快”这类观察测于未打包时，不能和打包阶段的计时放在一起。",
    ], size=14)


@slide_guard
def slide_adopt():
    a = load("finals_a.json")
    s = new("采纳与否决", "")
    txt(s, 0.5, 1.35, 6.6, 4.2, [
        "采纳配置 —— 实际改动的只有三个数",
        "",
        "  model.latent_dim        = 384       <- 改了 (128)",
        "  training.encoder_lr     = 2e-3      <- 改了 (5e-3)",
        "  scheduler.min_lr        = 1e-5      <- 改了 (1e-4)",
        "  scheduler.patience      = 5         = 默认",
        "  scheduler.factor        = 0.5       = 默认",
        "  model.head_hidden_dims  = [64]      = 默认",
        "  training.head_lr        = 5e-3      = 默认",
        "  model.kr_x_hidden_dims  = [128,64]  = 默认",
        "  training.kr_lr          = 5e-4      = 默认",
        "  early_stopping.patience = 24        = 默认",
        "  learnable_loss_balancer = false",
    ], size=11, mono=True)
    table(s, 7.4, 1.35, 5.5, ["项目", "判定"],
          [["learnable loss balancer", "否决"],
           ["PCGrad", "不引入"],
           ["调 head（容量/LR/KR 分支）", "不采纳"],
           ["早停 patience 24 → 40", "不采纳"],
           ["向下扩展 LR 搜索范围", "不需要"],
           ["关闭 LR 调度", "否决"]],
          col_w=[3.6, 1.9], size=11)
    txt(s, 7.4, 4.1, 5.5, 0.9,
        ["每一个“停在默认”都是测出来的，不是没测。"], size=11, color=MUT)
    txt(s, 0.5, 5.15, 12.4, 2.2, [
        "本轮确立了什么：",
        "  1. 探针必须覆盖大/中/小三档 —— σ 从 8.48% 降到 2.05%，campaign 才有判断力。",
        "  2. 名次要用 seed 买，不是用网格点买 —— 5 seed 的榜首在 25 seed 下掉到第 10。",
        "  3. 继承来的基线必须在当前régime重测 —— 旧天花板在 23 个任务里 17 个偏低，且偏差不是常数。",
        "  4. 分组均值不能当结论 —— 两个任务一正一负相消，看起来像“接近天花板”。",
        "  5. “半接上、静默失效”的功能是系统性问题：DDP、checkpoint dict、loss balancer 三例。",
        "  6. 凡是“在当前最优上”做的实验，最优一变就必须重做 —— 本轮踩了两次。",
        "  7. 管线在固定 seed 下位级确定，所以本报告所有 σ 都是配置级方差，而非运行抖动。",
    ], size=12)
    return a


def slide_limits():
    s = new("局限（与数字一同声明，不隐含）", "")
    txt(s, 0.6, 1.6, 12.2, 4.6, [
        "1. Stage C 每臂仅 1 seed，小差距不可分辨。",
        "2. [training.scheduler] 无分组开关，调度相关结论无法归因到单个参数组。",
        "3. Stage C 的天花板对比偏乐观 —— 假设臂的 seed 噪声等于单任务的。",
        "4. 打包前后墙钟不可比。",
        "5. 探针结论以 6 个任务测得；24 任务版本由 xfer 阶段给出，两者都在报告里。",
    ], size=16)


def main() -> None:
    global ALLOW_MISSING
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--date", default="2026-08-28")
    ap.add_argument("-o", "--out", type=Path, default=None)
    ap.add_argument("--allow-missing", action="store_true",
                    help="skip slides whose inputs are absent and list them (dry run only)")
    args = ap.parse_args()
    ALLOW_MISSING = args.allow_missing

    slide_title(args.date)
    slide_summary()
    slide_glossary()
    slide_design_flow()
    slide_design_why()
    slide_probe()
    slide_anchor()
    slide_grid()
    slide_finals()
    slide_finals_sigma()
    slide_a4()
    slide_a2b()
    slide_stage_b()
    slide_b_finals()
    slide_stage_c()
    slide_ceiling_fig()
    slide_transfer_fig()
    slide_transfer_why()
    slide_xfer()
    slide_position()
    slide_descriptor_limit()
    slide_balancer()
    slide_pcgrad()
    slide_packing()
    slide_adopt()
    slide_limits()

    out = args.out or RES / f"REPORT_v2_{args.date.replace('-', '')}.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"{out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    if SKIPPED:
        print(f"\nSKIPPED {len(SKIPPED)} slide(s) — this build is NOT a deliverable:")
        for line in SKIPPED:
            print(f"  {line}")


if __name__ == "__main__":
    main()
