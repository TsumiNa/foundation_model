#!/usr/bin/env python3
"""Build results/RECIPE_20260812.pptx — the compact companion deck to HYBRID_RECIPE.md:
the final training recipe (hybrid replay + optional consolidation), its validation, and
next-phase ops notes. 7 slides, meant for quick supplementary explanation.

Run: uv run --with python-pptx --with pillow python experiments/replay_epoch_sweep/build_recipe_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
AN = HERE / "analysis"
OUT = HERE / "results" / "RECIPE_20260812.pptx"

INK = RGBColor(0x1F, 0x29, 0x37)
MUT = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x00, 0x9E, 0x73)
AMBER = RGBColor(0xFD, 0xE6, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADGREY = RGBColor(0x4B, 0x55, 0x63)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def txt(slide, x, y, w, h, lines, size=14, color=INK, mono=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if mono:
            p.font.name = "Consolas"
    return box


def title_bar(slide, title, sub=None):
    t = txt(slide, 0.5, 0.25, 12.3, 0.6, [title], size=24)
    t.text_frame.paragraphs[0].font.bold = True
    if sub:
        txt(slide, 0.5, 0.85, 12.3, 0.4, [sub], size=12, color=MUT)


def pic_slide(title, sub, img, top=1.35, bottom=0.2):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title, sub)
    from PIL import Image

    with Image.open(img) as im:
        iw, ih = im.size
    max_w, max_h = Inches(12.9), prs.slide_height - Inches(top) - Inches(bottom)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(img), Emu(int((prs.slide_width - w) / 2)), Inches(top), Emu(w), Emu(h))
    return s


# 1 — title
s = prs.slides.add_slide(BLANK)
txt(s, 0.7, 2.3, 12, 1.6, ["Final Training Recipe",
                           "Hybrid replay during training + one optional consolidation pass"], size=32)
s.shapes[-1].text_frame.paragraphs[0].font.bold = True
txt(s, 0.7, 4.2, 12, 1.5, [
    "amount = max(1500 labels, 0.3 · N_task) per old task · resampled every epoch · patience 24, ≤150 epochs",
    "validated 2026-08-12 on the 24-task set (companion doc: HYBRID_RECIPE.md; full evidence: REPORT_20260802/09)",
], size=15, color=MUT)

# 2 — the recipe card
s = prs.slides.add_slide(BLANK)
title_bar(s, "The recipe — two steps, one optional",
          "full config: configs/hybrid_full24.toml · consolidation config: configs/joint_retrain_full24.toml")
txt(s, 0.6, 1.4, 6.3, 0.4, ["1 · Continual pretraining (mandatory)"], size=15)
s.shapes[-1].text_frame.paragraphs[0].font.bold = True
txt(s, 0.6, 1.9, 6.3, 3.4, [
    "[pretrain.replay]",
    "interval = 1",
    'resample = "epoch"      # redraw subset every epoch',
    "amount   = 0.30         # 30% of each old task",
    "per_task = { <every task with 0.3·N < 1500> = 1500 }",
    "",
    "[training]",
    "max_epochs = 150        # 100 loses only ~0.005",
    "early_stopping.patience = 24   # ≈ early stop off",
], size=12.5, mono=True)
txt(s, 0.6, 5.3, 6.3, 1.6, [
    "Rule for any task set:  amount_t = max(1500, 0.3·N_t)",
    "— global 0.3, floor 1500 on every task with N < 5000",
    "(engine clamps to N ⇒ small tasks auto-full-coverage)",
], size=12.5, color=MUT)
txt(s, 7.2, 1.4, 5.6, 0.4, ["2 · Consolidation (optional, ~1.5 h)"], size=15)
s.shapes[-1].text_frame.paragraphs[0].font.bold = True
txt(s, 7.2, 1.9, 5.6, 2.6, [
    "fm finetune \\",
    "  --config configs/joint_retrain_full24.toml \\",
    "  --checkpoint <out>/training/final_model.pt \\",
    "  --epochs 250 --output-dir <out>_joint",
    "",
    "# freeze_encoder = false, all 24 heads, full data",
    "# early-stops at ~76 epochs from a healthy model",
], size=12.5, mono=True)
txt(s, 7.2, 4.6, 5.6, 2.3, [
    "Take it when big-task (≥20k) performance matters:",
    "big-task deficit 0.031 → 0.022 (best of any arm).",
    "Skip it when small tasks are the priority (they give",
    "back a little: 0.008 → 0.016). It is polish, not rescue —",
    "from a no-replay model it only ever reaches 0.584.",
], size=12.5, color=MUT)

# 3 — validation results
s = prs.slides.add_slide(BLANK)
title_bar(s, "Validation — the two deliverables",
          "mean final R² over 23 R² tasks · deficit = single-task ceiling − final (group mean) · single seed, ±0.02 noise band")
VHEAD = ("arm", "mean R²", "big deficit", "mid deficit", "small deficit")
VROWS = [
    ("1 · hybrid replay", "0.652", "0.031", "0.012", "0.008", True),
    ("2 · hybrid + consolidation", "0.658", "0.022", "0.005", "0.016", True),
    ("best pure ratio (r0.3)", "0.652", "0.025", "0.008", "0.046", False),
    ("best pure fixed (n2500)", "0.663", "0.044", "−0.007", "0.002", False),
    ("no replay + retrain (control)", "0.584", "0.112", "0.061", "0.146", False),
]
vt = s.shapes.add_table(len(VROWS) + 1, 5, Inches(0.7), Inches(1.6), Inches(8.2), Inches(3.4)).table
for c, name in enumerate(VHEAD):
    cell = vt.cell(0, c)
    cell.text = name
    cell.fill.solid()
    cell.fill.fore_color.rgb = HEADGREY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
for r, (*row, hot) in enumerate(VROWS, start=1):
    for c, val in enumerate(row):
        cell = vt.cell(r, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = AMBER if hot else WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = hot or (c == 0)
        p.font.color.rgb = INK
txt(s, 9.2, 1.6, 3.7, 5.2, [
    "Result 1:",
    "minimax winner of all 12 replay",
    "settings — first arm to hold",
    "every size group ≤ 0.031",
    "",
    "Result 2:",
    "early stop @76 epochs (collapsed",
    "model needs 214) — replay already",
    "learned almost everything;",
    "minimax drops to 0.022",
    "",
    "The 2×2: replay during training",
    "sets the ceiling (0.584 vs 0.65+);",
    "end retraining adds only +0.006",
], size=12)

# 4-5 — evidence figures
pic_slide("Why the hybrid allocation — fixed-count starves BIG tasks, ratio starves SMALL ones",
          "deficit to the single-task ceiling vs labels actually replayed per old task, every arm on one axis; "
          "the black star (hybrid) is the only setting in/near the green zone in all three panels",
          AN / "replay_requirement_vs_size.png")
pic_slide("Why replay must happen DURING training",
          "left: per-step boxplots of test R² — without replay the step-24 median falls to −23 (symlog axis) · "
          "right: per-task distributions per retrain cap — converged (stop @214) with the whole box below the "
          "continual-replay reference",
          AN / "baseline_family.png")

# 6 — ops
s = prs.slides.add_slide(BLANK)
title_bar(s, "Ops notes for the next phase", "planning numbers from the validated runs (H200, 24-task set)")
txt(s, 0.6, 1.4, 6.2, 5.4, [
    "Wall-clock planning:",
    "· hybrid pretraining: ~68k replay labels/step → 21.6 h",
    "  (pure r0.3: 21.9 h · n2500: ~25 h with resume)",
    "· late steps cost 1–3 h each — kernel-regression replay",
    "  dominates; budget walltime accordingly",
    "· consolidation: ~1.5 h (early stop ~76 epochs)",
    "· A100 ≈ 1.3–1.5× these times",
], size=13)
txt(s, 7.0, 1.4, 5.8, 5.4, [
    "Recovery & correctness:",
    "· fm pretrain --resume is idempotent — resubmit the same",
    "  command after a TIMEOUT (only the in-flight step is lost)",
    "· resumed runs write PARTIAL metrics_table.csv — rebuild",
    "  from per-step JSONs (analysis/rebuild_metrics_from_stepdirs.py)",
    "· fm finetune has NO resume — give it walltime once",
    "· interval > 1 / no-replay usage needs the PR #36 fix",
    "  (in master since 921ffca); epoch resample is incompatible",
    "  with persistent_workers = true",
], size=13)

# 7 — caveats & checklist
s = prs.slides.add_slide(BLANK)
title_bar(s, "Caveats & next-phase checklist")
txt(s, 0.6, 1.4, 12.2, 5.4, [
    "Known limits of the evidence:",
    "· single seed (2025), single fixed task order — ±0.02 differences are noise; publish-grade numbers need ≥3 seeds",
    "· classification (material_type) is insensitive to replay settings (±0.005) — don't tune for it",
    "· open control: step-p24 (frozen subsets + long training) — the last cut separating coverage from training length",
    "",
    "Checklist when applying to a new task set:",
    "· compute N_t per task; set per_task = 1500 for every N_t < 5000 (rule: amount = max(1500, 0.3·N))",
    "· keep the m150 recipe: resample epoch + patience 24 + max_epochs 150",
    "· re-anchor any protocol sized under frozen-subset assumptions (task-scaling replay branches n1000/n1500)",
    "· decide consolidation by priority: big-task accuracy → yes; small-task sensitivity → skip",
], size=13.5)

OUT.parent.mkdir(exist_ok=True)
prs.save(OUT)
print(f"saved {OUT}")
