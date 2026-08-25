#!/usr/bin/env python3
"""Build results/REPORT_20260809.pptx — replay epoch-resampling sweep + budget variants +
ratio-replay family + no-replay/joint-retrain baselines.

Run: uv run --with python-pptx python experiments/replay_epoch_sweep/build_report_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
AN = HERE / "analysis"
OUT = HERE / "results" / "REPORT_20260809.pptx"

INK = RGBColor(0x1F, 0x29, 0x37)
MUT = RGBColor(0x6B, 0x72, 0x80)
ACC = RGBColor(0x00, 0x77, 0xBB)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def txt(slide, x, y, w, h, lines, size=14, bold_first=False, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
    return box


def title_bar(slide, title, sub=None):
    t = txt(slide, 0.5, 0.25, 12.3, 0.6, [title], size=24)
    t.text_frame.paragraphs[0].font.bold = True
    if sub:
        txt(slide, 0.5, 0.85, 12.3, 0.4, [sub], size=12, color=MUT)


def pic_slide(title, sub, img, top=1.35, bottom=0.2):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title, sub)
    from PIL import Image

    with Image.open(img) as im:
        iw, ih = im.size
    max_w, max_h = Inches(12.9), prs.slide_height - Inches(top) - Inches(bottom)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(img), Emu(int((prs.slide_width - w) / 2)), Inches(top), Emu(w), Emu(h))
    return s


# 1 — title
s = prs.slides.add_slide(BLANK)
txt(s, 0.7, 2.3, 12, 1.6, ["Replay with Per-Epoch Resampling",
                           "Fixed-count & ratio sweeps · training-budget variants · no-replay baselines"],
    size=32, bold_first=True)
txt(s, 0.7, 4.1, 12, 1.5, [
    "replay.resample = \"epoch\": redraw each old task's replay subset (count n or fraction r) every epoch",
    "34 runs · 7 arms · 3 machines | 2026-07-29 → 08-09",
    "step-p8 baseline (rikyu GB200, 2026-07) · epoch-p8 (ism 4×A100) · epoch-p24 / epoch-m150 / ratio-m150 (R-CCS H200)",
    "baselines: no-replay sequential (H200) → joint retrain at 4 epoch caps (H200 + 3×A100)",
], size=15, color=MUT)

# 2 — TL;DR
s = prs.slides.add_slide(BLANK)
title_bar(s, "TL;DR — five findings")
txt(s, 0.6, 1.3, 12.2, 5.8, [
    "1. Per-epoch resampling lifts retention at EVERY replay budget: mean final R² +0.022 … +0.126; the whole",
    "    saturation curve shifts up-left. Effective replay multiplier 3–5× (epoch-n200 ≈ step-n1000; epoch-n500 > step-n1500).",
    "",
    "2. The gain is coverage-shaped: largest at small n (peak +0.126 @ n200), shrinking monotonically to +0.022 @ n2500 —",
    "    exactly the N·(1−(1−n/N)^E) prediction. At n2500 the frozen subset already ≈ the full pool, and the two modes converge.",
    "",
    "3. Training length amplifies it: with patience 24 every step runs the full 100 epochs and the n-dependence nearly",
    "    flattens (0.59–0.65 over a 25× budget range). n100-p24 (0.592) ≈ n2500-step (0.600): 100 resampled labels ≈ 2500 frozen.",
    "",
    "4. The epoch budget saturates near ~100 across the WHOLE n range (m150 row complete): m150−p24 = −0.012…+0.024,",
    "    mean +0.005, no n-trend — even n100 does not keep climbing (0.580 vs 0.592). Coverage accumulation is capped by",
    "    the per-epoch replay gradient weight, not the epoch count.",
    "",
    "5. Early stopping was silently the binding constraint under resampling: fresh data each epoch delays the val-loss plateau",
    "    (mean 60–65 epochs/step vs 52–59 frozen at patience 8); at patience 24, 100% of steps hit the max_epochs cap.",
], size=13.5)

# 2b — TL;DR extension: ratio family + baselines
s = prs.slides.add_slide(BLANK)
title_bar(s, "Extension — ratio replay & the no-replay baselines (three more findings)",
          "ratio amount r ∈ {0.1, 0.2, 0.3, 0.5} of each old task's labels, m150 recipe · "
          "baseline: 24 tasks with NO replay, then ONE full-data joint retrain of all 24 heads")
txt(s, 0.6, 1.5, 12.2, 5.6, [
    "6. On the MEAN, ratio buys nothing new: put on the same cost axis (total labels replayed per step), ratio and",
    "    fixed-count land on one saturated curve (ratio 0.64–0.65 ≈ fixed n1000–n2500). Parameterization ≠ more performance.",
    "",
    "7. But allocation decides WHO pays: fixed-count starves big tasks (deficit plateaus at ~0.045 — what the previous",
    "    report called \"multi-task cost\" was half recoverable forgetting: ratio 0.3/0.5 cuts it to ~0.025). Ratio starves",
    "    small tasks symmetrically (r0.5 deficit 0.085 vs 0.002 at n2500, which replays 100% of their data).",
    "    ⇒ hybrid rule: amount = max(floor_n, r·N_task) — expressible today via replay.per_task overrides.",
    "",
    "8. Rehearsal timing is structural: without replay only 4% of learned tasks still beat a constant predictor by step 24",
    "    (mean R² −33). One full-data joint retrain at the end CONVERGES (early stop @214 of a 300-epoch cap, verified",
    "    caps 150/200/250/300) at 0.584 — below every continual-replay arm (0.639–0.663). Replaying during training,",
    "    even 100 labels/task, beats unlimited rehearsal after the damage is done.",
], size=13.5)

# 3 — design
s = prs.slides.add_slide(BLANK)
title_bar(s, "Design — one flag changed, everything else pinned",
          "canonical rikyu sweep configs (24 tasks, fixed order, seed 2025, batch 256, max 100 epochs + early stop)")
txt(s, 0.6, 1.4, 6.2, 5.6, [
    "Arms (fixed-count family only, n = 100…2500):",
    "· step-p8 — frozen subset per step (historical baseline)",
    "· epoch-p8 — only delta: --set replay.resample=\"epoch\"",
    "· epoch-p24 — + early_stopping.patience 8→24",
    "· epoch-m150 — + max_epochs 100→150 (4 heavy n)",
    "",
    "Execution:",
    "· epoch-p8: ism-gpu-a100, 4 workers, heavy-first FIFO",
    "· p24/m150: R-CCS ai-h200-brc, 1 GPU + 28 cores/job,",
    "  submissions trickled under the 72 node-hour quota,",
    "  walltime kills recovered by idempotent --resume",
], size=13, bold_first=False)
txt(s, 7.0, 1.4, 5.8, 5.6, [
    "Provenance:",
    "· every run: run_provenance.json (resolved config + git commit)",
    "· commits de711ed (ism) / 25a58b1 (R-CCS); configs untouched",
    "· metrics collected as results/mt_n*_{epoch,epoch_p24,epoch_m150}.csv",
    "",
    "Caveats (accepted by design):",
    "· single seed (2025), single fixed task order",
    "· cross-hardware arms (GB200 / A100 / H200)",
    "· mask-RNG protocol differs from the 2026-07 baseline",
    "  (statistically equivalent subsets, not bit-identical)",
    "· n2500-step ≈ epoch convergence acts as the negative control",
], size=13)

# 4-6 — figures
pic_slide("Headline — the saturation curve shifts up-left",
          "mean final test R² over 23 R² tasks vs replay n; every arm at every n above the frozen baseline",
          AN / "mean_final_compare.png")
pic_slide("Per-task saturation with both baselines — step vs epoch, 24 panels",
          "grey dashed = at-intro (step runs), red dashed = at-intro (epoch runs), green = single-task baseline; "
          "largest per-task gains: klat +0.166 @n200, tc +0.163 @n100; classification (material_type) unchanged",
          AN / "per_task_saturation_compare.png")
pic_slide("Forgetting stays event-driven — but small-n collapses are suppressed",
          "per-task trajectories across replay events; blues = frozen (light→dark = n100→n2500), reds = per-epoch resampling",
          AN / "replay_trajectories_compare.png")

# 6b — replay requirement vs task size (both baselines, all budgets)
s = pic_slide("Distance to the single-task ceiling — EVERY arm on one axis, by task size",
              "x = labels actually replayed per old task (group mean of min(amount, N), log) — ratio, hybrid and the "
              "end-of-run joint retrain join the four count arms in one frame; two dashed refs: at-intro of the "
              "p24/150 arms (anchors the green zone) vs the p8 arms (introductions early-stop lower)",
              AN / "replay_requirement_vs_size.png", top=1.5, bottom=1.0)
txt(s, 0.5, 6.7, 12.4, 0.7, [
    "big (≥20k): the count arms plateau at ~0.045; ratio r0.3/0.5 and the hybrid cut it to 0.025–0.031 — half the old "
    "\"multi-task cost\" was recoverable forgetting. small (≤1.2k): the mirror image — ratio starves them (r0.5: 0.085) "
    "while 100%-coverage settings (n≥1500, the hybrid floor) sit at ≤0.01. The grey diamond (no-replay → end retrain) "
    "is worst in every group; the hybrid star is the only setting in or at the green zone in all three panels.",
], size=11, color=MUT)

# 6c-6e — ratio family + baselines
pic_slide("Ratio replay joins the same cost curve — no free lunch on the mean",
          "mean final R² (23 tasks) vs TOTAL labels replayed per step; both families under the m150 recipe; "
          "references: converged end-of-run joint retrain (0.584) and the best frozen-subset arm (0.600)",
          AN / "ratio_cost_view.png")
pic_slide("Baseline family — collapse without replay; end-of-run rehearsal converges below every replay arm",
          "left: per-step boxplots of test R² over the learned tasks (symlog below −1; white diamond = mean) — without "
          "replay the whole distribution falls off the linear scale (step-24 medians: −23 vs 0.66) · right, one frame: "
          "retrain-from-collapse at 4 caps (converged @214, 250 ≡ 300, mean 0.584) vs the healthy trio — replay n1000 "
          "(0.643), hybrid (0.652), hybrid + consolidation (0.658, tightest box)",
          AN / "baseline_family.png")

# 6f — hybrid validation (the 2×2 closes)
from pptx.enum.text import PP_ALIGN

s = prs.slides.add_slide(BLANK)
title_bar(s, "Validation — the hybrid rule works; consolidation closes the 2×2",
          "hybrid = max(1500 labels, 0.3·N) per old task, m150 recipe (~68k labels/step, 21.6 h) · "
          "then ONE full-data joint retrain (cap 250, patience 24)")
VHEAD = ("arm", "mean R²", "big deficit", "mid deficit", "small deficit")
VROWS = [
    ("hybrid replay", "0.652", "0.031", "0.012", "0.008"),
    ("hybrid + joint retrain", "0.658", "0.022", "0.005", "0.016"),
    ("best pure ratio (r0.3)", "0.652", "0.025", "0.008", "0.046"),
    ("best pure fixed (n2500)", "0.663", "0.044", "−0.007", "0.002"),
    ("no replay + joint retrain", "0.584", "0.112", "0.061", "0.146"),
]
INKCLR = RGBColor(0x1F, 0x29, 0x37)
vt = s.shapes.add_table(len(VROWS) + 1, 5, Inches(0.7), Inches(1.7), Inches(7.6), Inches(3.2)).table
for c, name in enumerate(VHEAD):
    cell = vt.cell(0, c)
    cell.text = name
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x4B, 0x55, 0x63)
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
for r, row in enumerate(VROWS, start=1):
    for c, val in enumerate(row):
        cell = vt.cell(r, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xFD, 0xE6, 0x8A) if r == 2 else RGBColor(0xFF, 0xFF, 0xFF)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = (r == 2) or (c == 0)
        p.font.color.rgb = INKCLR
txt(s, 8.6, 1.7, 4.3, 5.0, [
    "Hybrid replay alone:",
    "first arm to hold EVERY size group ≤0.031 — the",
    "minimax winner among all 12 replay settings",
    "",
    "+ joint retrain (result 2):",
    "early-stops at 76 epochs (vs 214 from the collapsed",
    "model) — replay already learned almost everything;",
    "mean +0.006 (noise-level), but big-task deficit",
    "0.031 → 0.022: best big-task value of ANY arm.",
    "Small tasks give back a little (0.008 → 0.016).",
    "",
    "⇒ recipe: hybrid replay DURING training;",
    "a ~1.5 h consolidation pass is optional polish,",
    "not a rescue mechanism.",
], size=12)

# 7 — variants table

s = prs.slides.add_slide(BLANK)
title_bar(s, "Training-budget variants — patience 24 flattens n; epochs saturate near 100",
          "mean final test R² (23 tasks) · bold = best in row · amber pair = same score at 25× different replay budget")
HEAD = ("n", "step-p8", "epoch-p8", "epoch-p24", "epoch-m150")
VALS = [
    (100, 0.371, 0.475, 0.592, 0.580),
    (200, 0.420, 0.546, 0.606, 0.624),
    (500, 0.498, 0.582, 0.637, 0.632),
    (1000, 0.556, 0.612, 0.644, 0.643),
    (1500, 0.578, 0.621, 0.641, 0.660),
    (2000, 0.595, 0.632, 0.647, 0.642),
    (2500, 0.600, 0.622, 0.639, 0.663),
]
ARM_RGB = {1: RGBColor(0x00, 0x77, 0xBB), 2: RGBColor(0xCC, 0x33, 0x11),
           3: RGBColor(0x00, 0x9E, 0x73), 4: RGBColor(0xEE, 0x77, 0x33)}
WHITE, AMBER, HEADGREY = RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xFD, 0xE6, 0x8A), RGBColor(0x4B, 0x55, 0x63)
PAIR_CELLS = {(1, 3), (7, 1)}  # n100 × epoch-p24  ↔  n2500 × step-p8: 0.592 ≈ 0.600

tbl = s.shapes.add_table(len(VALS) + 1, 5, Inches(0.7), Inches(1.6), Inches(6.8), Inches(4.6)).table
for c, name in enumerate(HEAD):
    cell = tbl.cell(0, c)
    cell.text = name
    cell.fill.solid()
    cell.fill.fore_color.rgb = ARM_RGB.get(c, HEADGREY)
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
for r, row in enumerate(VALS, start=1):
    best = max(row[1:])
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = str(val) if c == 0 else f"{val:.3f}"
        cell.fill.solid()
        cell.fill.fore_color.rgb = AMBER if (r, c) in PAIR_CELLS else WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = (c == 0) or (c >= 1 and val == best) or (r, c) in PAIR_CELLS
        p.font.color.rgb = INK if c else HEADGREY

CATS = [
    ("Resampling (epoch-p8)", ARM_RGB[2], "+0.022…+0.126 vs frozen — positive at every n"),
    ("Full 100 epochs (p24)", ARM_RGB[3], "column goes flat: 0.59–0.65 across a 25× budget range"),
    ("Budget equivalence", None, "amber pair: 100 redrawn labels ≈ 2500 frozen (0.592 ≈ 0.600)"),
    ("150-epoch cap (m150)", ARM_RGB[4], "vs p24: mean +0.005, no n-trend — saturated at ~100 epochs"),
    ("Not yet run", None, "step-p24 control (frozen + long training)"),
]
y = 1.75
for label, color, desc in CATS:
    box = txt(s, 7.9, y, 5.0, 0.85, [label, desc], size=12.5)
    ps = box.text_frame.paragraphs
    ps[0].font.bold = True
    if color is not None:
        ps[0].font.color.rgb = color
    ps[1].font.size = Pt(11.5)
    ps[1].font.color.rgb = MUT
    y += 1.0

# 8 — next steps
s = prs.slides.add_slide(BLANK)
title_bar(s, "Practical guidance & next steps")
txt(s, 0.6, 1.4, 12.2, 5.4, [
    "Adopt now:",
    "· default pretrain.replay.resample = \"epoch\" — it never hurt at any n (worst case +0.022, no regression observed)",
    "· replay DURING training is non-negotiable: skipping it and jointly retraining at the end converges 0.055–0.08 below",
    "  every continual-replay arm — and hits big & small tasks hardest",
    "· size replay as a hybrid: amount = max(floor ≈ 500–1000, r·N_task) with r ≈ 0.1–0.3, via replay.per_task overrides —",
    "  fixes the fixed-count big-task deficit without the pure-ratio small-task starvation",
    "",
    "Cost note: ratio replay is wall-clock heavy — late steps at r=0.1 take ~3 h/step on H200 (kernel-regression replay",
    "dominates); r=0.5 ran 22.5 h for 24 steps. The hybrid floor keeps most of the gain at a fraction of the cost.",
    "",
    "Next:",
    "· step-p24 control (4 runs, ism) to split \"train longer\" from \"resample coverage\" at large n",
    "· multi-seed (≥3) confirmation of the headline deltas before publishing numbers",
    "· hybrid-rule validation run: max(500, 0.2·N) × epoch resampling vs n2500 and r0.3 at matched wall-clock",
    "",
    "Data & reproduction: experiments/replay_epoch_sweep/ (README, worker/launch/sbatch/collect scripts, analysis/*.py — in git);",
    "results mt_*.csv + joint_retrain_m*.json + figures + this deck under results/ & analysis/ (rsync policy, not in git).",
], size=13.5)

OUT.parent.mkdir(exist_ok=True)
prs.save(OUT)
print(f"saved {OUT}")
